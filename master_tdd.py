"""
╔══════════════════════════════════════════════════════════════════════════════╗
║          TRAINING DESIGN AGENT — Oracle University                          ║
║          Single-file Streamlit App · Groq LLaMA-3.3-70B                    ║
║                                                                              ║
║  SECTIONS (use Ctrl+F to jump):                                              ║
║   1. IMPORTS & CONSTANTS                                                     ║
║   2. PROMPTS  (No-Labs / Yes-Labs system prompts)                            ║
║   3. SESSION STATE                                                           ║
║   4. THEME & CSS                                                             ║
║   5. FILE UTILITIES  (extraction, OCR, copyright scan)                       ║
║   6. URL CRAWLER                                                             ║
║   7. LLM CLIENT  (Groq call, prompt builder, quality check)                 ║
║   8. EXPORT  (Oracle-branded DOCX + PDF)                                    ║
║   9. SCREEN 1 — Course Information & Target Audience                        ║
║  10. SCREEN 2 — Source Content                                               ║
║  11. SCREEN 3 — Generate, Audit & Export                                    ║
║  12. MAIN ROUTER                                                             ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

# ══════════════════════════════════════════════════════════════════════════════
# 1. IMPORTS & CONSTANTS
# ══════════════════════════════════════════════════════════════════════════════

import io
import os
import re
import html as html_lib
import logging
import subprocess
import tempfile
import urllib.request
import urllib.parse
import urllib.error
from datetime import datetime
from html.parser import HTMLParser
from pathlib import Path
from typing import Optional

import streamlit as st

logging.basicConfig(level=logging.WARNING)
logger = logging.getLogger(__name__)

# ── Oracle brand colours ──────────────────────────────────────────────────────
ORACLE_RED      = "#C74634"
ORACLE_DARK     = "#3A3A3A"
ORACLE_GREY     = "#F5F5F5"
ORACLE_BORDER   = "#E0E0E0"
ORACLE_RED_HEX  = "C74634"
ORACLE_DARK_HEX = "3A3A3A"
ORACLE_GREY_HEX = "F0F0F0"

# ── LLM settings ──────────────────────────────────────────────────────────────
GROQ_MODEL  = "llama-3.3-70b-versatile"
MAX_TOKENS  = 8000

# ── File ingestion limits ─────────────────────────────────────────────────────
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".txt", ".csv"}
BLOCKED_EXTENSIONS = {".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv", ".aac", ".flac", ".ogg"}

# ── Copyright scanner keyword patterns ───────────────────────────────────────
COPYRIGHT_PATTERNS = {
    "copyright": [r"©", r"\bcopyright\b", r"all rights reserved"],
    "trademark": [r"™", r"®", r"\btrademark\b"],
    "restricted": [r"\bconfidential\b", r"\bproprietary\b", r"do not distribute"],
    "licensed":  [r"\blicensed\b", r"license agreement", r"unauthorized use prohibited"],
}

# ── URL crawler limits ────────────────────────────────────────────────────────
MAX_LINKS_PER_LEVEL = 15
MAX_CRAWL_CHARS     = 40_000
CRAWL_DEPTH         = 2
CRAWL_TIMEOUT       = 10

# ── Required sections for QA check ───────────────────────────────────────────
REQUIRED_SECTIONS_BASE = [
    "Course Overview", "Course Title", "Product Area", "Training Need",
    "Target Audience", "Learning Objectives", "Benefits to Learner",
    "Course Description", "Assumptions", "Implementation Readiness",
    "GTM Messaging", "Course Coverage", "Case Study",
    "QA Checklist", "Traceability Map",
]

STEP_LABELS = {
    1: "📋  Course Information",
    2: "📂  Source Content",
    3: "🚀  Generate & Export",
}


# ══════════════════════════════════════════════════════════════════════════════
# 2. PROMPTS
# ══════════════════════════════════════════════════════════════════════════════

NO_LABS_SYSTEM_PROMPT = """
You are an expert Oracle University Instructional Designer AI agent. Your task is to create a complete, learner-centric Training Design Document aligned to modern instructional design standards. The design must be practical, implementation-ready, and mapped directly to real job performance.

PROCESS & GUARDRAILS (REQUIRED):
- Use only the provided Product Documentation / Source Content for product capabilities; do not invent features, UI flows, terminology, or configurations. If information is missing, capture it under "Assumptions & Open Questions."
- Keep all content strictly within the specified product scope and target job roles. Explicitly exclude out-of-scope topics.
- If any critical input seems missing or unclear, note it under Assumptions & Open Questions — do NOT fabricate information.

INSTRUCTIONAL DESIGN STANDARDS:
- Map each job task to specific required skills. Ensure every module supports measurable on-the-job performance.
- Apply the 80-20 principle. Include a brief 80/20 rationale (1-2 bullets).
- Write measurable, performance-based learning objectives using Bloom's Taxonomy action verbs.
- Assign Bloom's levels for each topic (Remember, Understand, Apply, Analyze, Evaluate, Create).

COURSE STRUCTURE:
- Learning Progression: Foundational → Intermediate → Advanced
- Content Mix per module: Concepts → Demonstrations → Scenarios
- Each Topic (one video) duration target: 3-7 minutes. Include estimated total course seat time.
- Each module must include at least: 1 Concept topic, 1 Demonstration topic, 1 Scenario/Job-based exercise.

DESIGN PRINCIPLES (apply to every piece of content):
- Start with the user. Teach tasks, not just tools. Simplify with purpose. Show, don't just tell.

PERSONA PROFILES: For each persona include: name, role, top 5 responsibilities, top 3 pain points, learning preferences.

LEARNING OUTCOMES: 5-8 course-level outcomes + one per module. All must be SMART using Bloom's action verbs.

SKILL CHECKS: 2 scenario-based questions per module with 4 options (A/B/C/D), plausible distractors, one correct answer tied to a learning outcome.

GTM MESSAGING must cover all 5: what the product is, what makes it stand out, who it's for, what problems it solves, what learners take away.

TRACEABILITY: Tag every factual claim with [FILE: filename], [URL: https://...], or [INPUT]. At the end, add a TRACEABILITY MAP section as a markdown table: | Source Tag | Reference Detail | Document Section |

OUTPUT — follow this EXACT structure (do not omit any section):

---
## Course Overview
**Course Title:**
**Product Area:**
**Training Need:**
**Target Audience:**

**Learning Objectives:**
1.
2.
3.
4.
5.

**80/20 Prioritization Rationale:**
- 
- 

**Benefits to Learner:**
**Course Description:**

**Assumptions & Open Questions:**
- Assumptions:
- Open Questions:

---
## Persona Profiles
| Persona Name | Role | Top 5 Responsibilities | Top 3 Pain Points | Learning Preferences |
|---|---|---|---|---|

---
## Implementation Readiness
**Prerequisites (learner):**
**Prerequisites (access):**
**Required tools/materials:**
**Accessibility & delivery:**
**Assessment plan:**

---
## GTM Messaging
**What the product is:**
**What makes it stand out:**
**Who the course is for:**
**What business problems it helps solve:**
**What learners will take away:**

---
## Course Coverage
| MODULE # | MODULE TITLE | MODULE LEARNING OBJECTIVE | TOPIC LIST | MATCHING SCENARIO EXERCISE |
|---|---|---|---|---|

(For TOPIC LIST: Topic Title | Bloom Level | Measurable Topic Objective | Est. Video Duration (min))
(For MATCHING SCENARIO EXERCISE: Scenario Title | Scenario Type | Success Criteria)

---
## Skill Checks
(2 scenario-based questions per module, 4 options each, correct answer identified)

---
## Case Study
**Goal:**
**Scenario:**
**Requirement:**
**Steps to Implement:**
**Expected Outcome:**

---
## QA Checklist
- [ ] Every job task maps to at least one skill and module/topic
- [ ] 80/20 prioritization rationale included
- [ ] Bloom level assigned for every topic; objectives use measurable action verbs
- [ ] Balanced mix achieved per module (Concept + Demo + Scenario)
- [ ] No out-of-scope content; gaps captured under Assumptions & Open Questions
- [ ] All learning outcomes are SMART
- [ ] GTM messaging covers all 5 required elements

---
## TRACEABILITY MAP
| Source Tag | Reference Detail | Document Section |
|---|---|---|

Write in a professional, clear, learner-focused tone. Fully populate every section. No placeholders.
"""

YES_LABS_SYSTEM_PROMPT = """
You are an expert Oracle University Instructional Designer AI agent. Your task is to create a complete, learner-centric Training Design Document aligned to modern instructional design standards. The design must be practical, implementation-ready, and mapped directly to real job performance.

PROCESS & GUARDRAILS (REQUIRED):
- Use only the provided Product Documentation / Source Content for product capabilities; do not invent features, UI flows, terminology, or configurations. If information is missing, capture it under "Assumptions & Open Questions."
- Keep all content strictly within the specified product scope and target job roles. Explicitly exclude out-of-scope topics.
- If any critical input seems missing or unclear, note it under Assumptions & Open Questions — do NOT fabricate information.

INSTRUCTIONAL DESIGN STANDARDS:
- Map each job task to specific required skills. Ensure every module supports measurable on-the-job performance.
- Apply the 80-20 principle. Include a brief 80/20 rationale (1-2 bullets).
- Write measurable, performance-based learning objectives using Bloom's Taxonomy action verbs.
- Assign Bloom's levels for each topic (Remember, Understand, Apply, Analyze, Evaluate, Create).

COURSE STRUCTURE:
- Learning Progression: Foundational → Intermediate → Advanced
- Content Mix per module: Concepts → Demonstrations → Labs → Scenarios
- Each Topic (one video) duration: 3-7 minutes. Provide duration per topic AND per lab. Include estimated total course seat time.
- Each module must include at least: 1 Concept topic, 1 Demonstration topic, 1 Hands-on Lab, 1 Scenario/Job-based exercise.

DESIGN PRINCIPLES (apply to every piece of content):
- Start with the user. Teach tasks, not just tools. Simplify with purpose. Show, don't just tell.

PERSONA PROFILES: For each persona include: name, role, top 5 responsibilities, top 3 pain points, learning preferences.

LEARNING OUTCOMES: 5-8 course-level outcomes + one per module. All must be SMART using Bloom's action verbs.

HANDS-ON LABS: Simulate real-world Oracle tasks. Include both guided and unguided types. Provide success criteria (1-2 bullets) per lab.

SKILL CHECKS: 2 scenario-based questions per module with 4 options (A/B/C/D), plausible distractors, one correct answer tied to a learning outcome.

GTM MESSAGING must cover all 5: what the product is, what makes it stand out, who it's for, what problems it solves, what learners take away.

TRACEABILITY: Tag every factual claim with [FILE: filename], [URL: https://...], or [INPUT]. At the end, add a TRACEABILITY MAP section as a markdown table: | Source Tag | Reference Detail | Document Section |

OUTPUT — follow this EXACT structure (do not omit any section):

---
## Course Overview
**Course Title:**
**Product Area:**
**Training Need:**
**Target Audience:**

**Learning Objectives:**
1.
2.
3.
4.
5.

**80/20 Prioritization Rationale:**
- 
- 

**Benefits to Learner:**
**Course Description:**

**Assumptions & Open Questions:**
- Assumptions:
- Open Questions:

---
## Persona Profiles
| Persona Name | Role | Top 5 Responsibilities | Top 3 Pain Points | Learning Preferences |
|---|---|---|---|---|

---
## Implementation Readiness
**Prerequisites (learner):**
**Prerequisites (access):**
**Required tools/materials:**
**Accessibility & delivery:**
**Assessment plan:**

---
## GTM Messaging
**What the product is:**
**What makes it stand out:**
**Who the course is for:**
**What business problems it helps solve:**
**What learners will take away:**

---
## Course Coverage
| MODULE # | MODULE TITLE | MODULE LEARNING OBJECTIVE | TOPIC LIST | MATCHING HANDS-ON LAB | HANDS-ON LAB DURATION (MINUTES) |
|---|---|---|---|---|---|

(For TOPIC LIST: Topic Title | Bloom Level | Measurable Topic Objective | Est. Video Duration (min))
(For MATCHING HANDS-ON LAB: Lab Title | Lab Type (guided/unguided) | Success Criteria)

---
## Skill Checks
(2 scenario-based questions per module, 4 options each, correct answer identified)

---
## Case Study
**Goal:**
**Scenario:**
**Requirement:**
**Steps to Implement:**
**Expected Outcome:**

---
## QA Checklist
- [ ] Every job task maps to at least one skill and module/topic
- [ ] 80/20 prioritization rationale included
- [ ] Bloom level assigned for every topic; objectives use measurable action verbs
- [ ] Balanced mix achieved per module (Concept + Demo + Lab + Scenario)
- [ ] No out-of-scope content; gaps captured under Assumptions & Open Questions
- [ ] All learning outcomes are SMART
- [ ] GTM messaging covers all 5 required elements
- [ ] Each module includes at least 1 lab with success criteria

---
## TRACEABILITY MAP
| Source Tag | Reference Detail | Document Section |
|---|---|---|

Write in a professional, clear, learner-focused tone. Fully populate every section. No placeholders.
"""


# ══════════════════════════════════════════════════════════════════════════════
# 3. SESSION STATE
# ══════════════════════════════════════════════════════════════════════════════

def init_session():
    defaults = {
        "step": 1,
        # Screen 1
        "course_title": "", "product_name": "", "context": "",
        "target_job_roles": "", "job_task_analysis": "",
        "course_type": "eLearning", "labs_required": False,
        "audience_level": "Beginner", "prerequisite_knowledge": "",
        # Screen 2
        "doc_links": [""], "uploaded_files_meta": [],
        "additional_text": "", "golden_standard_text": "",
        "crawled_content": {},
        # Generation
        "generated_doc": "", "generation_done": False,
        "traceability_rows": [], "source_counts": {},
        # Feedback
        "user_feedback": "", "regeneration_count": 0,
        # Export
        "docx_bytes": None, "pdf_bytes": None,
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v


# ══════════════════════════════════════════════════════════════════════════════
# 4. THEME & CSS
# ══════════════════════════════════════════════════════════════════════════════

def apply_theme():
    st.markdown(f"""
    <style>
    html, body, [class*="css"] {{ font-family: 'Arial', sans-serif; }}

    .oracle-header {{
        background: {ORACLE_RED}; color: white; padding: 14px 28px;
        border-radius: 6px; margin-bottom: 24px;
    }}
    .oracle-header h1 {{ margin: 0; font-size: 22px; font-weight: 700; }}
    .oracle-header p  {{ margin: 2px 0 0; font-size: 13px; opacity: 0.88; }}

    .section-card {{
        background: white; border: 1px solid {ORACLE_BORDER}; border-radius: 8px;
        padding: 22px 26px; margin-bottom: 18px;
        box-shadow: 0 1px 4px rgba(0,0,0,0.06);
    }}
    .section-card h3 {{
        color: {ORACLE_RED}; margin-top: 0; font-size: 15px;
        border-bottom: 2px solid {ORACLE_RED}; padding-bottom: 7px; margin-bottom: 14px;
    }}

    .stButton > button {{
        background: {ORACLE_RED} !important; color: white !important;
        border: none !important; border-radius: 5px !important;
        font-weight: 600 !important; padding: 8px 22px !important;
    }}
    .stButton > button:hover {{ background: #a83828 !important; }}

    section[data-testid="stSidebar"] {{ background: {ORACLE_DARK}; }}
    section[data-testid="stSidebar"] * {{ color: white !important; }}

    .copyright-warning {{
        background: #fff3cd; border: 1px solid #ffc107;
        border-left: 5px solid #e6a200; border-radius: 4px; padding: 12px 16px; margin: 8px 0;
    }}
    .audit-panel {{
        background: {ORACLE_GREY}; border: 1px solid {ORACLE_BORDER};
        border-radius: 6px; padding: 16px 20px;
    }}
    .doc-output {{
        background: white; border: 1px solid {ORACLE_BORDER}; border-radius: 6px;
        padding: 22px; max-height: 680px; overflow-y: auto;
        font-size: 13.5px; line-height: 1.7;
    }}
    .tooltip-text {{ font-size: 11.5px; color: #777; margin-top: -5px; margin-bottom: 8px; }}
    .stProgress > div > div > div > div {{ background-color: {ORACLE_RED}; }}
    </style>

    <div class="oracle-header">
        <h1>🎓 Training Design Agent</h1>
        <p>Oracle University · AI-Powered Instructional Design</p>
    </div>
    """, unsafe_allow_html=True)


def render_sidebar():
    step = st.session_state.get("step", 1)
    with st.sidebar:
        st.markdown("## Progress")
        for s, label in STEP_LABELS.items():
            icon = "✅" if s < step else ("🔵" if s == step else "⬜")
            st.markdown(f"**{icon} {label}**" if s == step else f"{icon} {label}")
        st.markdown("---")
        st.caption("Powered by Groq LLaMA-3.3-70B\n\nGenerates learner-centric training design documents aligned to Oracle instructional design standards.")
        if step > 1 and st.button("⬅ Back to Step 1"):
            st.session_state["step"] = 1; st.rerun()
        if step > 2 and st.button("⬅ Back to Step 2"):
            st.session_state["step"] = 2; st.rerun()


def tip(text):
    st.markdown(f'<p class="tooltip-text">ℹ️ {text}</p>', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# 5. FILE UTILITIES
# ══════════════════════════════════════════════════════════════════════════════

def validate_extension(filename: str) -> tuple[bool, str]:
    ext = Path(filename).suffix.lower()
    if ext in BLOCKED_EXTENSIONS:
        return False, f"Audio/video files ({ext}) are not supported."
    if ext not in ALLOWED_EXTENSIONS:
        return False, f"Unsupported file type: {ext}. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
    return True, ""


def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":            return _extract_pdf(file_bytes)
        elif ext in (".docx",".doc"):return _extract_docx(file_bytes)
        elif ext in (".pptx",".ppt"):return _extract_pptx(file_bytes)
        elif ext in (".xlsx",".xls"):return _extract_xlsx(file_bytes)
        elif ext in (".txt",".csv"): return file_bytes.decode("utf-8", errors="replace")
        else: return ""
    except Exception as e:
        return f"[Error extracting {filename}: {e}]"


def _extract_pdf(file_bytes: bytes) -> str:
    parts = []
    # Try pdfplumber first
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t: parts.append(t)
    except ImportError:
        pass
    if parts: return "\n".join(parts)

    # Try PyMuPDF
    try:
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc: parts.append(page.get_text())
        if any(t.strip() for t in parts): return "\n".join(parts)
    except Exception:
        pass

    # OCR fallback
    return _ocr_pdf(file_bytes)


def _ocr_pdf(file_bytes: bytes) -> str:
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
        images = convert_from_bytes(file_bytes, dpi=200)
        return "\n".join(pytesseract.image_to_string(img) for img in images)
    except ImportError:
        return ("[OCR unavailable — pytesseract/pdf2image not installed. "
                "This appears to be a scanned/image-based PDF. "
                "Please provide a text-based PDF or paste content manually.]")
    except Exception as e:
        return f"[OCR failed: {e}]"


def _extract_docx(file_bytes: bytes) -> str:
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(c.text for c in row.cells))
        return "\n".join(parts)
    except Exception as e: return f"[DOCX error: {e}]"


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"\n--- Slide {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append("\t".join(c.text for c in row.cells))
        return "\n".join(parts)
    except Exception as e: return f"[PPTX error: {e}]"


def _extract_xlsx(file_bytes: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"\n--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                parts.append("\t".join(str(c) if c is not None else "" for c in row))
        return "\n".join(parts)
    except Exception as e: return f"[XLSX error: {e}]"


def scan_copyright(text: str) -> dict:
    sample = text[:3000].lower()
    flags = {}
    for category, patterns in COPYRIGHT_PATTERNS.items():
        matched = [p for p in patterns if re.search(p, sample, re.IGNORECASE)]
        if matched: flags[category] = matched
    return flags


def copyright_warning_msg(flags: dict) -> str:
    bullets = "\n".join(f"• **{cat.capitalize()}** indicators detected" for cat in flags)
    return ("⚠️ **This document may contain protected content.**\n\n" + bullets +
            "\n\nPlease confirm you have the rights to use this for training purposes.")


# ══════════════════════════════════════════════════════════════════════════════
# 6. URL CRAWLER
# ══════════════════════════════════════════════════════════════════════════════

class _LinkExtractor(HTMLParser):
    def __init__(self, base_url):
        super().__init__()
        self.base_url = base_url
        self.links = []

    def handle_starttag(self, tag, attrs):
        if tag != "a": return
        href = dict(attrs).get("href", "")
        if not href or href.startswith(("#","javascript:","mailto:")): return
        full = urllib.parse.urljoin(self.base_url, href)
        if urllib.parse.urlparse(full).netloc == urllib.parse.urlparse(self.base_url).netloc:
            self.links.append(full)


def _fetch_url(url: str) -> Optional[str]:
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "OracleTrainingAgent/1.0"})
        with urllib.request.urlopen(req, timeout=CRAWL_TIMEOUT) as r:
            return r.read().decode(r.headers.get_content_charset("utf-8"), errors="replace")
    except Exception as e:
        logger.warning(f"Fetch failed {url}: {e}")
        return None


def _html_to_text(html: str) -> str:
    text = re.sub(r"<script[^>]*>.*?</script>", " ", html, flags=re.DOTALL|re.IGNORECASE)
    text = re.sub(r"<style[^>]*>.*?</style>",  " ", text,  flags=re.DOTALL|re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"\s{3,}", "\n\n", text)
    return text.strip()


def crawl_url(base_url: str) -> tuple[str, list[str]]:
    if not base_url.startswith(("http://","https://")):
        base_url = "https://" + base_url
    visited, to_visit = set(), [(base_url, 0)]
    all_parts, crawled, total = [], [], 0

    while to_visit and total < MAX_CRAWL_CHARS:
        url, depth = to_visit.pop(0)
        if url in visited: continue
        visited.add(url)
        html = _fetch_url(url)
        if not html: continue
        text = _html_to_text(html)
        snippet = text[:MAX_CRAWL_CHARS - total]
        all_parts.append(f"\n\n=== SOURCE: {url} ===\n{snippet}")
        total += len(snippet)
        crawled.append(url)
        if depth < CRAWL_DEPTH:
            ex = _LinkExtractor(url)
            ex.feed(html)
            seen = set()
            for l in ex.links:
                if l not in visited and l not in seen:
                    seen.add(l)
                    to_visit.append((l, depth+1))
                    if len(seen) >= MAX_LINKS_PER_LEVEL: break

    return "\n".join(all_parts), crawled


# ══════════════════════════════════════════════════════════════════════════════
# 7. LLM CLIENT
# ══════════════════════════════════════════════════════════════════════════════

def get_groq_client():
    try:
        from groq import Groq
        return Groq(api_key=st.secrets["GROQ_API_KEY"])
    except KeyError:
        st.error("❌ **GROQ_API_KEY not found in Streamlit secrets.**\n\n"
                 "Go to **App Settings → Secrets** and add:\n```\nGROQ_API_KEY = \"your_key_here\"\n```")
        st.stop()
    except ImportError:
        st.error("❌ `groq` package not installed. Add `groq` to requirements.txt.")
        st.stop()


def build_user_prompt() -> str:
    ss = st.session_state

    # Assemble source content
    source_parts = []
    for meta in ss.get("uploaded_files_meta", []):
        if meta.get("confirmed") or not meta.get("flagged"):
            source_parts.append(f"[FILE: {meta['name']}]\n{meta['text'][:8000]}")

    for url, text in ss.get("crawled_content", {}).items():
        source_parts.append(f"[URL: {url}]\n{text[:6000]}")

    if ss.get("additional_text","").strip():
        source_parts.append(f"[INPUT]\n{ss['additional_text']}")

    source_block = "\n\n---\n\n".join(source_parts) if source_parts else "(No source content provided.)"

    golden = ""
    if ss.get("golden_standard_text","").strip():
        golden = ("\n\n## GOLDEN STANDARD REFERENCE\n"
                  "Match the tone, depth, and writing style of the following approved reference document. "
                  "Do NOT deviate from the required template structure — only adapt style.\n\n"
                  + ss["golden_standard_text"][:3000])

    feedback = ""
    if ss.get("user_feedback","").strip() and ss.get("regeneration_count",0) > 0:
        feedback = ("\n\n## REVISION INSTRUCTIONS\n"
                    "Incorporate the following feedback precisely. Do not rewrite sections not mentioned:\n\n"
                    + ss["user_feedback"])

    return f"""
# INPUTS

- **Course Title:** {ss.get('course_title','')} [INPUT]
- **Product Name:** {ss.get('product_name','')} [INPUT]
- **Context:** {ss.get('context','')} [INPUT]
- **Target Job Roles:** {ss.get('target_job_roles','')} [INPUT]
- **Job Task Analysis:** {ss.get('job_task_analysis','')} [INPUT]
- **Course Type:** {ss.get('course_type','eLearning')} [INPUT]
- **Audience Experience Level:** {ss.get('audience_level','Beginner')} [INPUT]
- **Prerequisite Knowledge:** {ss.get('prerequisite_knowledge','None specified')} [INPUT]

# PRODUCT DOCUMENTATION / SOURCE CONTENT

{source_block}
{golden}
{feedback}

Now generate the complete Training Design Document following ALL template sections exactly.
""".strip()


def generate_doc(progress_cb=None) -> str:
    client = get_groq_client()
    labs   = st.session_state.get("labs_required", False)
    system = YES_LABS_SYSTEM_PROMPT if labs else NO_LABS_SYSTEM_PROMPT

    if progress_cb: progress_cb(0.08, "Sending inputs to Groq LLaMA-3.3-70B...")

    full_text = ""
    try:
        stream = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {"role": "system", "content": system},
                {"role": "user",   "content": build_user_prompt()},
            ],
            max_tokens=MAX_TOKENS,
            temperature=0.3,
            stream=True,
        )
        chunk_n = 0
        for chunk in stream:
            delta = chunk.choices[0].delta.content or ""
            full_text += delta
            chunk_n += 1
            if progress_cb and chunk_n % 40 == 0:
                pct = min(0.08 + (len(full_text) / (MAX_TOKENS * 4)) * 0.85, 0.93)
                progress_cb(pct, "Generating design document…")
    except Exception as e:
        raise RuntimeError(f"Groq API call failed: {e}") from e

    if progress_cb: progress_cb(0.96, "Finalising…")
    return full_text


def quality_check(doc_text: str) -> dict:
    required = REQUIRED_SECTIONS_BASE + (["Hands-On Lab"] if st.session_state.get("labs_required") else [])
    present  = [s for s in required if re.search(re.escape(s), doc_text, re.IGNORECASE)]
    missing  = [s for s in required if s not in present]
    return {"present": present, "missing": missing, "pass": not missing}


def extract_traceability(doc_text: str) -> tuple[list[dict], dict]:
    rows, counts = [], {}
    tag_pat  = re.compile(r"\[(FILE:[^\]]+|URL:[^\]]+|INPUT)\]")
    current  = "General"
    for line in doc_text.split("\n"):
        hm = re.match(r"^#{1,3}\s+(.+)", line)
        if hm: current = hm.group(1).strip()
        for m in tag_pat.finditer(line):
            tag = m.group(1)
            rows.append({"Source Tag": tag, "Document Section": current, "Context": line[:80]})
            counts[tag] = counts.get(tag, 0) + 1
    return rows, counts


# ══════════════════════════════════════════════════════════════════════════════
# 8. EXPORT  (Oracle-branded DOCX + PDF)
# ══════════════════════════════════════════════════════════════════════════════

def build_docx(md_text: str) -> bytes:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        raise ImportError("python-docx not installed. Add `python-docx` to requirements.txt.")

    C_RED   = RGBColor(0xC7, 0x46, 0x34)
    C_DARK  = RGBColor(0x3A, 0x3A, 0x3A)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def cell_bg(cell, hex_col):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"),"clear"); shd.set(qn("w:color"),"auto"); shd.set(qn("w:fill"), hex_col)
        tcPr.append(shd)

    def hrule(doc, col=ORACLE_GREY_HEX):
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr(); pBdr = OxmlElement("w:pBdr")
        bot = OxmlElement("w:bottom")
        bot.set(qn("w:val"),"single"); bot.set(qn("w:sz"),"6")
        bot.set(qn("w:space"),"1"); bot.set(qn("w:color"), col)
        pBdr.append(bot); pPr.append(pBdr)
        p.paragraph_format.space_after = Pt(3)

    def inline_run(para, text):
        parts = re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", text)
        for part in parts:
            if part.startswith("**") and part.endswith("**"):
                r = para.add_run(part[2:-2]); r.font.bold = True
            elif part.startswith("*") and part.endswith("*"):
                r = para.add_run(part[1:-1]); r.font.italic = True
            else:
                r = para.add_run(part)
            para.runs[-1].font.name = "Arial"
            para.runs[-1].font.size = Pt(10.5)

    def add_md_table(doc, table_lines):
        rows = []
        for line in table_lines:
            line = line.strip()
            if not line.startswith("|"): continue
            if re.match(r"^\|[-| :]+\|$", line): continue
            cells = [c.strip() for c in line.split("|") if c.strip() != ""]
            if cells: rows.append(cells)
        if not rows: return
        cols = max(len(r) for r in rows)
        tbl = doc.add_table(rows=len(rows), cols=cols)
        tbl.style = "Table Grid"
        tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
        for ri, row in enumerate(rows):
            for ci in range(cols):
                txt = row[ci] if ci < len(row) else ""
                cell = tbl.cell(ri, ci)
                cell.vertical_alignment = WD_ALIGN_VERTICAL.TOP
                p = cell.paragraphs[0]; p.clear()
                run = p.add_run(txt); run.font.name = "Arial"; run.font.size = Pt(9)
                if ri == 0:
                    run.font.bold = True; run.font.color.rgb = C_WHITE
                    cell_bg(cell, ORACLE_RED_HEX)
                elif ri % 2 == 0:
                    cell_bg(cell, ORACLE_GREY_HEX)
        doc.add_paragraph()

    doc = Document()
    sec = doc.sections[0]
    sec.page_width  = int(8.5 * 914400)
    sec.page_height = int(11  * 914400)
    for attr in ("left_margin","right_margin","top_margin","bottom_margin"):
        setattr(sec, attr, Inches(1))
    doc.styles["Normal"].font.name = "Arial"
    doc.styles["Normal"].font.size = Pt(10.5)

    # Cover page
    ss = st.session_state
    cover_tbl = doc.add_table(rows=1, cols=1)
    cover_tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
    cc = cover_tbl.cell(0,0); cell_bg(cc, ORACLE_RED_HEX)
    cp = cc.paragraphs[0]
    cr = cp.add_run("  Oracle University")
    cr.font.color.rgb = C_WHITE; cr.font.bold = True
    cr.font.size = Pt(13); cr.font.name = "Arial"
    cp.paragraph_format.space_before = Pt(8); cp.paragraph_format.space_after = Pt(8)

    doc.add_paragraph()
    tp = doc.add_paragraph()
    tr = tp.add_run(ss.get("course_title","Training Design Document"))
    tr.font.size = Pt(24); tr.font.bold = True
    tr.font.color.rgb = C_RED; tr.font.name = "Arial"

    doc.add_paragraph()
    for label, val in [
        ("Product Area",    ss.get("product_name","")),
        ("Course Type",     ss.get("course_type","")),
        ("Target Audience", ss.get("target_job_roles","")),
        ("Audience Level",  ss.get("audience_level","")),
        ("Labs Included",   "Yes" if ss.get("labs_required") else "No"),
        ("Document Date",   datetime.today().strftime("%B %d, %Y")),
    ]:
        p = doc.add_paragraph()
        rl = p.add_run(f"{label}: "); rl.font.bold = True
        rl.font.size = Pt(11); rl.font.name = "Arial"; rl.font.color.rgb = C_DARK
        rv = p.add_run(str(val))
        rv.font.size = Pt(11); rv.font.name = "Arial"
        p.paragraph_format.space_after = Pt(3)

    doc.add_page_break()

    # Parse markdown
    lines = md_text.split("\n")
    i, table_lines, in_table = 0, [], False

    while i < len(lines):
        raw = lines[i].rstrip()

        if raw.startswith("|"):
            table_lines.append(raw); in_table = True; i += 1; continue

        if in_table:
            add_md_table(doc, table_lines)
            table_lines = []; in_table = False

        if raw.startswith("## "):
            hrule(doc, ORACLE_RED_HEX)
            p = doc.add_paragraph()
            r = p.add_run(raw[3:].strip())
            r.font.size = Pt(14); r.font.bold = True
            r.font.color.rgb = C_RED; r.font.name = "Arial"
            p.paragraph_format.space_before = Pt(14); p.paragraph_format.space_after = Pt(5)
        elif raw.startswith("### "):
            p = doc.add_paragraph()
            r = p.add_run(raw[4:].strip())
            r.font.size = Pt(12); r.font.bold = True
            r.font.color.rgb = C_DARK; r.font.name = "Arial"
            p.paragraph_format.space_before = Pt(10); p.paragraph_format.space_after = Pt(3)
        elif raw.startswith("# "):
            p = doc.add_paragraph()
            r = p.add_run(raw[2:].strip())
            r.font.size = Pt(18); r.font.bold = True
            r.font.color.rgb = C_RED; r.font.name = "Arial"
            p.paragraph_format.space_before = Pt(18); p.paragraph_format.space_after = Pt(6)
        elif raw.startswith("---"):
            hrule(doc)
        elif raw.startswith("- [ ]") or raw.startswith("- [x]"):
            p = doc.add_paragraph(style="List Bullet")
            tick = "☑ " if raw.startswith("- [x]") else "☐ "
            r = p.add_run(tick + raw[5:].strip())
            r.font.name = "Arial"; r.font.size = Pt(10)
        elif raw.startswith("- ") or raw.startswith("* "):
            p = doc.add_paragraph(style="List Bullet")
            inline_run(p, raw[2:].strip())
        elif re.match(r"^\d+\.\s", raw):
            p = doc.add_paragraph(style="List Number")
            inline_run(p, re.sub(r"^\d+\.\s","",raw).strip())
        elif raw.strip() == "":
            doc.add_paragraph().paragraph_format.space_after = Pt(2)
        else:
            p = doc.add_paragraph()
            inline_run(p, raw)
            p.paragraph_format.space_after = Pt(4)

        i += 1

    if table_lines: add_md_table(doc, table_lines)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_pdf(docx_bytes: bytes) -> Optional[bytes]:
    with tempfile.TemporaryDirectory() as tmp:
        dp = os.path.join(tmp, "doc.docx")
        pp = os.path.join(tmp, "doc.pdf")
        with open(dp,"wb") as f: f.write(docx_bytes)
        try:
            r = subprocess.run(
                ["libreoffice","--headless","--convert-to","pdf","--outdir",tmp,dp],
                capture_output=True, timeout=60)
            if r.returncode == 0 and os.path.exists(pp):
                return open(pp,"rb").read()
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
    return None


def safe_filename(name: str) -> str:
    return re.sub(r"[^\w\s-]","",name).strip().replace(" ","_")[:60] or "design_doc"


# ══════════════════════════════════════════════════════════════════════════════
# 9. SCREEN 1 — Course Information & Target Audience
# ══════════════════════════════════════════════════════════════════════════════

def screen1():
    st.markdown("### Step 1 of 3 — Course Information & Target Audience")
    st.caption("Fill in all required fields. The more detail you provide, the more accurate and tailored the output.")

    # Course Information
    st.markdown('<div class="section-card"><h3>📘 Course Information</h3>', unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    with c1:
        st.session_state["course_title"] = st.text_input(
            "Course Title *", value=st.session_state["course_title"],
            placeholder="e.g. Oracle HCM Cloud: Absence Management for HR Admins")
    with c2:
        st.session_state["product_name"] = st.text_input(
            "Product Name / Area *", value=st.session_state["product_name"],
            placeholder="e.g. Oracle HCM Cloud — Absence Management")

    st.session_state["context"] = st.text_area(
        "Context *", value=st.session_state["context"], height=130,
        placeholder=("Provide background and purpose: What is new? Why does this training matter? "
                     "What product feature or update does it cover?"))
    tip("Be as specific as possible — include release notes, new features, and business context.")
    st.markdown("</div>", unsafe_allow_html=True)

    # Target Audience
    st.markdown('<div class="section-card"><h3>👥 Target Audience</h3>', unsafe_allow_html=True)
    c3, c4 = st.columns(2)
    with c3:
        st.session_state["target_job_roles"] = st.text_input(
            "Target Job Roles *", value=st.session_state["target_job_roles"],
            placeholder="e.g. HR Administrator, Benefits Specialist, Payroll Manager")
    with c4:
        levels = ["Beginner", "Intermediate", "Advanced"]
        st.session_state["audience_level"] = st.selectbox(
            "Audience Experience Level", options=levels,
            index=levels.index(st.session_state["audience_level"]))

    st.session_state["job_task_analysis"] = st.text_area(
        "Job Task Analysis (Focus Areas) *", value=st.session_state["job_task_analysis"], height=150,
        placeholder=("Summarize key tasks per role:\n"
                     "HR Administrator: Configure absence plans, manage requests, run reports...\n"
                     "Benefits Specialist: Enroll employees, manage life events..."))
    tip("Describe 4–6 real job tasks per role — what they do daily in the system.")

    st.session_state["prerequisite_knowledge"] = st.text_input(
        "Prerequisite Knowledge or Skills (Optional)", value=st.session_state["prerequisite_knowledge"],
        placeholder="e.g. Basic Oracle Cloud navigation, HCM Foundations course")
    st.markdown("</div>", unsafe_allow_html=True)

    # Course Settings
    st.markdown('<div class="section-card"><h3>⚙️ Course Settings</h3>', unsafe_allow_html=True)
    c5, c6 = st.columns(2)
    with c5:
        types = ["eLearning","Instructor-Led Training (ILT)","Virtual Instructor-Led Training (vILT)","Blended (eLearning + ILT)"]
        ct = st.session_state["course_type"]
        if ct not in types: ct = "eLearning"
        st.session_state["course_type"] = st.selectbox("Course Type", options=types, index=types.index(ct))
    with c6:
        st.session_state["labs_required"] = st.toggle(
            "🔬 Labs Required For This Course?",
            value=st.session_state["labs_required"],
            help="ON = Hands-on Labs in every module (Yes-Labs template). OFF = Scenarios only (No-Labs template).")
        if st.session_state["labs_required"]:
            st.success("✅ Labs ON — each module will include a hands-on lab.")
        else:
            st.info("ℹ️ Labs OFF — each module will include scenario-based exercises.")
    st.markdown("</div>", unsafe_allow_html=True)

    # Navigation
    st.markdown("---")
    _, col_next = st.columns([3,1])
    with col_next:
        if st.button("Next: Source Content →", use_container_width=True):
            errors = []
            for field, label in [("course_title","Course Title"), ("product_name","Product Name"),
                                  ("context","Context"), ("target_job_roles","Target Job Roles"),
                                  ("job_task_analysis","Job Task Analysis")]:
                if not st.session_state.get(field,"").strip():
                    errors.append(f"❌ {label} is required.")
            if errors:
                for e in errors: st.error(e)
            else:
                st.session_state["step"] = 2; st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# 10. SCREEN 2 — Source Content
# ══════════════════════════════════════════════════════════════════════════════

def screen2():
    st.markdown("### Step 2 of 3 — Source Content")
    st.caption("Upload reference materials and provide documentation links. All content is extracted and analyzed before generation.")

    # Documentation Links
    st.markdown('<div class="section-card"><h3>🔗 Documentation Links</h3>', unsafe_allow_html=True)
    tip("Enter a base URL — the crawler follows child links 2 levels deep automatically.")
    links = st.session_state.get("doc_links", [""])
    updated = []
    for i, link in enumerate(links):
        cl, cd = st.columns([6,1])
        with cl:
            val = st.text_input(f"URL {i+1}", value=link, key=f"url_{i}",
                                placeholder="https://docs.oracle.com/en/cloud/saas/...")
            updated.append(val)
        with cd:
            st.write("")
            if len(links) > 1 and st.button("✕", key=f"del_url_{i}"):
                links.pop(i); st.session_state["doc_links"] = links; st.rerun()
    st.session_state["doc_links"] = updated
    if st.button("➕ Add another URL"):
        st.session_state["doc_links"].append(""); st.rerun()
    st.markdown("</div>", unsafe_allow_html=True)

    # File Uploads
    st.markdown('<div class="section-card"><h3>📂 Upload Supporting Documents</h3>', unsafe_allow_html=True)
    tip("Accepted: PDF, DOCX, PPTX, XLSX, TXT, CSV. Scanned PDFs will be OCR-processed. No audio/video.")

    uploaded = st.file_uploader(
        "Upload documents (multiple allowed)", accept_multiple_files=True,
        type=["pdf","docx","doc","pptx","ppt","xlsx","xls","txt","csv"])

    existing_meta  = st.session_state.get("uploaded_files_meta", [])
    existing_names = {m["name"] for m in existing_meta}

    if uploaded:
        for uf in uploaded:
            if uf.name in existing_names: continue
            ok, reason = validate_extension(uf.name)
            if not ok:
                st.error(f"❌ **{uf.name}**: {reason}"); continue
            with st.spinner(f"Processing {uf.name}…"):
                fb   = uf.read()
                text = extract_text(fb, uf.name)
            flags = scan_copyright(text)
            existing_meta.append({
                "name": uf.name, "text": text,
                "flagged": bool(flags), "confirmed": not bool(flags), "flags": flags,
            })
            existing_names.add(uf.name)
        st.session_state["uploaded_files_meta"] = existing_meta

    remove_idx = []
    for idx, meta in enumerate(st.session_state.get("uploaded_files_meta", [])):
        with st.expander(f"📄 {meta['name']}", expanded=meta["flagged"]):
            st.caption(f"Extracted: {len(meta['text']):,} characters")
            if meta["flagged"]:
                st.markdown(f'<div class="copyright-warning">{copyright_warning_msg(meta["flags"])}</div>',
                            unsafe_allow_html=True)
                confirmed = st.checkbox(
                    "✅ I confirm I have the rights to use this document for training purposes.",
                    key=f"confirm_{idx}_{meta['name']}", value=meta.get("confirmed", False))
                st.session_state["uploaded_files_meta"][idx]["confirmed"] = confirmed
                if not confirmed:
                    st.warning("⚠️ This file will be **excluded** from generation until you confirm usage rights.")
            else:
                st.success("✅ No copyright flags detected. Ready for analysis.")
            if st.button("🗑 Remove", key=f"rm_{idx}_{meta['name']}"):
                remove_idx.append(idx)

    if remove_idx:
        st.session_state["uploaded_files_meta"] = [
            m for i,m in enumerate(st.session_state["uploaded_files_meta"]) if i not in remove_idx]
        st.rerun()
    st.markdown("</div>", unsafe_allow_html=True)

    # Additional Text
    st.markdown('<div class="section-card"><h3>📝 Additional Content</h3>', unsafe_allow_html=True)
    st.session_state["additional_text"] = st.text_area(
        "Paste additional technical information, notes or content (optional)",
        value=st.session_state.get("additional_text",""), height=140,
        placeholder="Paste release notes, feature descriptions, SME notes, process flows…")
    st.markdown("</div>", unsafe_allow_html=True)

    # Golden Standard
    st.markdown('<div class="section-card"><h3>⭐ Golden Standard Reference (Optional)</h3>', unsafe_allow_html=True)
    tip("Upload an approved reference design doc. AI will match its tone/depth — template structure is always preserved.")
    gf = st.file_uploader("Upload reference doc (DOCX or PDF)", type=["pdf","docx"], key="golden_upload")
    if gf:
        with st.spinner("Extracting reference tone…"):
            gt = extract_text(gf.read(), gf.name)
        st.session_state["golden_standard_text"] = gt
        st.success(f"✅ Golden Standard loaded: {gf.name} ({len(gt):,} chars)")
    elif st.session_state.get("golden_standard_text",""):
        st.info("ℹ️ A Golden Standard reference is already loaded.")
        if st.button("🗑 Clear Golden Standard"):
            st.session_state["golden_standard_text"] = ""; st.rerun()
    st.markdown("</div>", unsafe_allow_html=True)

    # Navigation
    st.markdown("---")
    cb, _, cn = st.columns([1,2,1])
    with cb:
        if st.button("⬅ Back to Step 1"): st.session_state["step"] = 1; st.rerun()
    with cn:
        if st.button("Generate Document →", use_container_width=True):
            # Crawl URLs
            valid_urls = [u.strip() for u in st.session_state.get("doc_links",[]) if u.strip()]
            if valid_urls:
                crawled = {}
                ph = st.empty()
                for url in valid_urls:
                    ph.info(f"🔍 Crawling: {url}")
                    try:
                        text, visited = crawl_url(url)
                        if text.strip():
                            crawled[url] = text
                            ph.success(f"✅ Crawled {len(visited)} page(s) from {url}")
                    except Exception as e:
                        ph.warning(f"⚠️ Could not crawl {url}: {e}")
                st.session_state["crawled_content"] = crawled
                ph.empty()
            st.session_state["step"] = 3; st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# 11. SCREEN 3 — Generate, Audit & Export
# ══════════════════════════════════════════════════════════════════════════════

def _md_preview(md: str) -> str:
    """Lightweight markdown → HTML for the preview panel."""
    lines = md.split("\n")
    out, table_buf, in_table = [], [], False

    def flush_table():
        rows = []
        for tl in table_buf:
            if re.match(r"^\|[-| :]+\|$", tl.strip()): continue
            cells = [c.strip() for c in tl.split("|") if c.strip()]
            if cells: rows.append(cells)
        if not rows: return ""
        cols = max(len(r) for r in rows)
        h = '<table style="border-collapse:collapse;width:100%;font-size:12px;margin:10px 0">'
        for ri, row in enumerate(rows):
            h += "<tr>"
            for ci in range(cols):
                v = html_lib.escape(row[ci]) if ci < len(row) else ""
                if ri == 0:
                    h += f'<th style="background:#C74634;color:white;padding:5px 9px;text-align:left;border:1px solid #ddd">{v}</th>'
                else:
                    bg = "#f8f8f8" if ri%2==0 else "white"
                    h += f'<td style="padding:4px 9px;border:1px solid #ddd;background:{bg}">{v}</td>'
            h += "</tr>"
        return h + "</table>"

    def inline(text):
        text = html_lib.escape(text)
        text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
        text = re.sub(r"\*(.+?)\*",     r"<em>\1</em>", text)
        text = re.sub(r"`(.+?)`",       r"<code>\1</code>", text)
        return text

    for line in lines:
        raw = line.rstrip()
        if raw.startswith("|"):
            table_buf.append(raw); in_table = True; continue
        if in_table:
            out.append(flush_table()); table_buf = []; in_table = False

        if raw.startswith("## "):
            txt = html_lib.escape(raw[3:])
            out.append(f'<h2 style="color:#C74634;font-size:15px;margin:16px 0 5px;border-bottom:2px solid #C74634;padding-bottom:3px">{txt}</h2>')
        elif raw.startswith("### "):
            out.append(f'<h3 style="color:#3A3A3A;font-size:13px;margin:11px 0 3px">{html_lib.escape(raw[4:])}</h3>')
        elif raw.startswith("# "):
            out.append(f'<h1 style="color:#C74634;font-size:19px;margin:18px 0 7px">{html_lib.escape(raw[2:])}</h1>')
        elif raw.startswith("---"):
            out.append('<hr style="border:none;border-top:1px solid #e0e0e0;margin:8px 0"/>')
        elif raw.startswith("- [ ]") or raw.startswith("- [x]"):
            tick = "☑" if raw.startswith("- [x]") else "☐"
            out.append(f'<p style="margin:2px 0;padding-left:14px">{tick} {inline(raw[5:].strip())}</p>')
        elif raw.startswith("- ") or raw.startswith("* "):
            out.append(f'<p style="margin:2px 0;padding-left:14px">• {inline(raw[2:].strip())}</p>')
        elif re.match(r"^\d+\.\s", raw):
            out.append(f'<p style="margin:2px 0;padding-left:14px">{inline(re.sub(r"^\d+\.\s","",raw).strip())}</p>')
        elif raw.strip() == "":
            out.append("<br/>")
        else:
            out.append(f'<p style="margin:3px 0;line-height:1.65">{inline(raw)}</p>')

    if in_table: out.append(flush_table())
    return "\n".join(out)


def screen3():
    st.markdown("### Step 3 of 3 — Generate & Export")

    # ── Auto-generate on first arrival ────────────────────────────────────────
    if not st.session_state.get("generation_done"):
        prog  = st.progress(0)
        status = st.empty()
        err    = st.empty()

        def cb(pct, msg):
            prog.progress(min(pct,1.0)); status.markdown(f"**{msg}**")

        try:
            cb(0.05, "🔍 Analyzing inputs and source content…")
            doc_text = generate_doc(cb)
            st.session_state["generated_doc"]  = doc_text
            st.session_state["generation_done"] = True
            rows, counts = extract_traceability(doc_text)
            st.session_state["traceability_rows"] = rows
            st.session_state["source_counts"]     = counts
            prog.progress(1.0)
            status.success("✅ Design document generated successfully!")
        except Exception as e:
            err.error(f"❌ Generation failed: {e}")
            if st.button("⬅ Go Back and Retry"):
                st.session_state["step"] = 2; st.rerun()
            return

    doc_text = st.session_state.get("generated_doc","")

    # ── Quality Check ──────────────────────────────────────────────────────────
    qc = quality_check(doc_text)
    with st.expander("🔍 Quality Check — Section Completeness", expanded=not qc["pass"]):
        if qc["pass"]:
            st.success(f"✅ All {len(qc['present'])} required sections are present.")
        else:
            st.warning(f"⚠️ {len(qc['missing'])} required section(s) may be missing.")
            c1, c2 = st.columns(2)
            with c1:
                st.markdown("**✅ Present:**")
                for s in qc["present"]: st.markdown(f"- {s}")
            with c2:
                st.markdown("**❌ Missing:**")
                for s in qc["missing"]: st.markdown(f"- {s}")

    # ── Document Preview ───────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("#### 📄 Generated Training Design Document")
    st.markdown(f'<div class="doc-output">{_md_preview(doc_text)}</div>', unsafe_allow_html=True)
    with st.expander("📋 View Raw Markdown"):
        st.code(doc_text, language="markdown")

    # ── Reliability Audit ──────────────────────────────────────────────────────
    st.markdown("---")
    rows   = st.session_state.get("traceability_rows",[])
    counts = st.session_state.get("source_counts",{})

    with st.expander("📊 Reliability Audit & Traceability Map", expanded=False):
        st.markdown('<div class="audit-panel">', unsafe_allow_html=True)
        st.markdown("#### Reliability Audit Summary")
        st.metric("Total Source Tags Found", len(rows))
        if counts:
            st.markdown("**Source → Usage Count:**")
            for src, n in sorted(counts.items(), key=lambda x: -x[1]):
                st.markdown(f"- `{src}` — **{n}** reference(s)")
        else:
            st.info("No inline source tags detected. Ensure source content was provided.")

        st.markdown("---")
        st.markdown("#### Source → Section Map")
        if rows:
            import pandas as pd
            df = pd.DataFrame(rows)[["Source Tag","Document Section","Context"]]
            st.dataframe(df, use_container_width=True, hide_index=True)
        else:
            st.info("No traceability data to display.")
        st.markdown("</div>", unsafe_allow_html=True)

    # ── Feedback & Regeneration ────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("#### ✏️ Provide Feedback & Regenerate")
    st.caption("Describe the changes you'd like. The AI incorporates feedback without rewriting correct sections.")

    feedback = st.text_area("Your feedback", value=st.session_state.get("user_feedback",""), height=110,
        placeholder=("e.g. 'Add a third module covering Advanced Reporting. "
                     "Make personas more specific to Finance roles. "
                     "Expand the Case Study with more detail.'"))
    st.session_state["user_feedback"] = feedback

    _, col_regen = st.columns([3,1])
    with col_regen:
        rc = st.session_state.get("regeneration_count",0)
        lbl = f"🔄 Regenerate (#{rc+1})" if rc > 0 else "🔄 Regenerate with Feedback"
        if st.button(lbl, use_container_width=True, disabled=not feedback.strip()):
            st.session_state["generation_done"] = False
            st.session_state["regeneration_count"] = rc + 1
            st.session_state["docx_bytes"] = None
            st.session_state["pdf_bytes"]  = None
            st.rerun()

    # ── Export ─────────────────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("#### ⬇️ Export")
    fname = safe_filename(st.session_state.get("course_title","design_doc"))
    cd, cp = st.columns(2)

    with cd:
        if st.button("📄 Prepare DOCX"):
            with st.spinner("Building Oracle-branded DOCX…"):
                try:
                    st.session_state["docx_bytes"] = build_docx(doc_text)
                except Exception as e:
                    st.error(f"DOCX error: {e}")
        if st.session_state.get("docx_bytes"):
            st.download_button("⬇️ Download DOCX", data=st.session_state["docx_bytes"],
                file_name=fname+".docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    with cp:
        if st.button("📕 Prepare PDF"):
            if not st.session_state.get("docx_bytes"):
                with st.spinner("Building DOCX first…"):
                    try: st.session_state["docx_bytes"] = build_docx(doc_text)
                    except Exception as e: st.error(f"DOCX error: {e}"); st.stop()
            with st.spinner("Converting to PDF via LibreOffice…"):
                pdf = build_pdf(st.session_state["docx_bytes"])
                if pdf: st.session_state["pdf_bytes"] = pdf
                else:
                    st.warning("⚠️ PDF conversion requires LibreOffice on the server. "
                               "Download the DOCX and convert locally, or add `libreoffice` to packages.txt.")
        if st.session_state.get("pdf_bytes"):
            st.download_button("⬇️ Download PDF", data=st.session_state["pdf_bytes"],
                file_name=fname+".pdf", mime="application/pdf")

    st.markdown("---")
    if st.button("🔁 Start a New Design Document"):
        for k in list(st.session_state.keys()): del st.session_state[k]
        st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# 12. MAIN ROUTER
# ══════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Training Design Agent | Oracle University",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded",
)

init_session()
apply_theme()
render_sidebar()

step = st.session_state.get("step", 1)
if   step == 1: screen1()
elif step == 2: screen2()
elif step == 3: screen3()
